import os
from pptx import Presentation
import logging

# إعداد السجلات
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

def process_pptx(input_path, output_path, placeholders, image_path=None, qr_path=None):
    """
    معالجة ملف العرض التقديمي واستبدال النصوص والصور (شعار gg_gg وباركود qr_qr)
    """
    try:
        logger.debug(f"بدء معالجة ملف العرض التقديمي")
        
        # التأكد من وجود ملف الإدخال
        if not os.path.exists(input_path):
            # محاولة البحث عن القالب في المجلد الثابت
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            template_path = os.path.join(base_dir, 'static', 'inputt.pptx')
            if os.path.exists(template_path):
                input_path = template_path
            else:
                raise FileNotFoundError(f"ملف القالب غير موجود في: {input_path}")

        # التأكد من وجود مجلد الإخراج
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # تحميل العرض التقديمي
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # استبدال الصور (الشعار والباركود)
                if shape.has_text_frame:
                    text = shape.text
                    target_image = None
                    if "gg_gg" in text and image_path and os.path.exists(image_path):
                        target_image = image_path
                    elif "qr_qr" in text and qr_path and os.path.exists(qr_path):
                        target_image = qr_path
                    
                    if target_image:
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                        slide.shapes.add_picture(target_image, left, top, width=width, height=height)
                        continue

                    # استبدال النصوص
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key, value in placeholders.items():
                                if key in run.text:
                                    run.text = run.text.replace(key, str(value))
        
        prs.save(output_path)
        return True
    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}")
        return False
